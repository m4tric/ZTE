from flask import Flask, render_template, request, send_file
from docx import Document
from docx.shared import Inches
from pptx import Presentation
import os
from datetime import datetime
import werkzeug

app = Flask(__name__)

# 🏠 Página principal con dos botones (Contratos y PO)
@app.route("/")
def home():
    return render_template("home.html")

# 📊 Formulario para Contrato (.doc)
@app.route("/contrato", methods=["GET", "POST"])
def contrato():
    return render_template("form.html")

# 📊 Formulario para PO (PowerPoint)
@app.route("/po", methods=["GET", "POST"])
def po():
    if request.method == "POST":
        user_data = request.form.to_dict()
        image_paths = {}

        # Manejo de imágenes
        for key in ["imagen_de_referencia1", "correo_requerimiento", "coti1", "coti2", "coti3", "cotifinal"]:
            if key in request.files:
                image = request.files[key]
                if image.filename:
                    safe_filename = werkzeug.utils.secure_filename(image.filename)
                    image_path = os.path.join(UPLOAD_FOLDER, safe_filename)
                    image.save(image_path)
                    image_paths[key] = image_path
                    print(f"✅ Imagen guardada en: {image_path}")
                else:
                    print(f"❌ No se recibió una imagen para {key}.")
            else:
                print(f"⚠️ No se encontró el campo '{key}' en request.files. **Revisar si el formulario tiene enctype='multipart/form-data'**")

        pptx_path = fill_pptx("PO_TEMPLATE.pptx", user_data, image_paths)
        download_filename = f"PO_{user_data.get('num_pro', 'SIN_NUM')}_{user_data.get('nom_pro', 'SIN_NOMBRE')}.pptx"
        return send_file(pptx_path, as_attachment=True, download_name=download_filename)
    return render_template("form_po.html")

# 🔹 Función mejorada para insertar imagen en PowerPoint
def fill_pptx(template_path, data, image_paths):
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"❌ Error al abrir el archivo de plantilla: {e}")
        return None

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in data.items():
                            if f"<{key}>" in run.text:
                                run.text = run.text.replace(f"<{key}>", value)
                                print(f"✅ Reemplazado {key} con valor: {value}")

    # 📌 Insertar imágenes en marcadores
    markers = {
        "<imagen_de_referencia1>": image_paths.get("imagen_de_referencia1"),
        "<correo_requerimiento>": image_paths.get("correo_requerimiento"),
        "<coti1>": image_paths.get("coti1"),
        "<coti2>": image_paths.get("coti2"),
        "<coti3>": image_paths.get("coti3"),
        "<cotifinal>": image_paths.get("cotifinal"),
    }

    for slide in prs.slides:
        for marker, image_path in markers.items():
            if image_path and os.path.exists(image_path):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            if marker in paragraph.text:
                                paragraph.text = ""
                                try:
                                    # ✅ Insertar imagen en el slide
                                    left = Inches(2)
                                    top = Inches(2)
                                    width = Inches(4)
                                    height = Inches(3)
                                    slide.shapes.add_picture(image_path, left, top, width, height)
                                    print(f"✅ Imagen insertada correctamente para el marcador {marker}.")
                                except Exception as e:
                                    print(f"❌ Error al insertar la imagen para {marker}: {e}")

    output_path = "generated_presentation.pptx"
    prs.save(output_path)
    return output_path

UPLOAD_FOLDER = "uploads"
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# 📌 Función para formatear fechas correctamente
def format_date(date_str, lang):
    if not date_str:
        return ""  
    try:
        date_obj = datetime.strptime(date_str, "%Y-%m-%d")
        if lang == "es":
            return date_obj.strftime("%d de %B de %Y").replace("March", "marzo").replace("April", "abril") \
                .replace("May", "mayo").replace("June", "junio").replace("July", "julio") \
                .replace("August", "agosto").replace("September", "septiembre") \
                .replace("October", "octubre").replace("November", "noviembre").replace("December", "diciembre")
        elif lang == "en":
            day = date_obj.day
            suffix = "th" if 11 <= day <= 13 else {1: "st", 2: "nd", 3: "rd"}.get(day % 10, "th")
            return date_obj.strftime(f"%B {day}{suffix}, %Y")
    except ValueError:
        return date_str  

# 📌 Función para reemplazar texto y agregar la imagen en Word
def fill_contract(template_path, user_data, image_path):
    doc = Document(template_path)

    print("\n📄 📩 DATOS QUE SE INSERTARÁN EN EL CONTRATO:")
    for key, value in user_data.items():
        print(f"{key}: {value}")  

    for paragraph in doc.paragraphs:
        for key, value in user_data.items():
            paragraph.text = paragraph.text.replace(f"<{key}>", value).replace(f"< {key} >", value)

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for key, value in user_data.items():
                    cell.text = cell.text.replace(f"<{key}>", value).replace(f"< {key} >", value)

    # 📌 Insertar la imagen en marcador <ANEXO_IMAGEN>
    if image_path and os.path.exists(image_path):
        print(f"🖼️ Insertando imagen desde: {image_path}")
        for paragraph in doc.paragraphs:
            if "<ANEXO_IMAGEN>" in paragraph.text:
                paragraph.clear()
                run = paragraph.add_run()
                run.add_picture(image_path, width=Inches(3))
                print("✅ Imagen insertada en un párrafo.")

        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if "<ANEXO_IMAGEN>" in cell.text:
                        cell.text = ""
                        run = cell.paragraphs[0].add_run()
                        run.add_picture(image_path, width=Inches(3))
                        print("✅ Imagen insertada en una celda de tabla.")
    else:
        print("⚠️ No se encontró la imagen o la ruta es incorrecta.")

    output_path = "Contrato_Generado.docx"
    doc.save(output_path)
    print("✅ Contrato generado correctamente:", output_path)
    return output_path

@app.route("/", methods=["GET", "POST"])
def index():
    fields = {
        "AGREEMENT_NO": "Número de requerimiento",
        "FECHA_FIRMA": "Fecha de firma",
        "FECHA_FIRMA_EN": "Fecha de firma en inglés",
        "FECHA_FINALIZACION": "Fecha de finalización del contrato",
        "FECHA_FINALIZACION_EN": "Fecha de finalización en inglés",
        "PROVEEDOR_NOMBRE": "Nombre del proveedor",
        "PROVEEDOR_DIRECCION": "Dirección del proveedor",
        "REPRESENTANTE_NOMBRE": "Nombre del representante legal del proveedor",
        "PROYECTO_NOMBRE": "Nombre del proyecto",
        "PROYECTO_NOMBRE_EN": "Nombre del proyecto en inglés",
        "PROYECTO_RESUMEN": "Resumen del proyecto",
        "PROYECTO_RESUMEN_EN": "Resumen del proyecto en inglés",
        "PROYECTO_MONTO": "Monto del proyecto",
        "PROYECTO_UBICACION": "Ubicación del proyecto",
        "BANCO_CUENTA_NOMBRE": "Nombre de la cuenta bancaria",
        "BANCO_CUENTA_NUMERO": "Número de cuenta",
        "BANCO_NOMBRE": "Nombre del banco",
        "BANCO_CLABE": "CLABE bancaria",
        "BANCO_DIRECCION": "Dirección del banco",
        "ANEXO_SERVICIOS": "Descripción de los servicios en español",
        "ANEXO_SERVICIOS_EN": "Descripción de los servicios en inglés",
    }

    if request.method == "POST":
        form_data = {key: request.form[key] for key in fields.keys() if key in request.form}

        image_path = None
        if "anexo_imagen" in request.files:
            image = request.files["anexo_imagen"]
            if image.filename:
                safe_filename = werkzeug.utils.secure_filename(image.filename)
                image_path = os.path.join(UPLOAD_FOLDER, safe_filename)
                image.save(image_path)
                print(f"✅ Imagen guardada en: {image_path}")
            else:
                print("❌ No se recibió una imagen.")
        else:
            print("⚠️ No se encontró el campo 'anexo_imagen' en request.files. **Revisar si el formulario tiene enctype='multipart/form-data'**")

        for date_key in ["FECHA_FIRMA", "FECHA_FIRMA_EN", "FECHA_FINALIZACION", "FECHA_FINALIZACION_EN"]:
            if date_key in form_data:
                lang = "es" if "EN" not in date_key else "en"
                form_data[date_key] = format_date(form_data[date_key], lang)

        contract_path = fill_contract("contract_template.docx", form_data, image_path)

        return send_file(contract_path, as_attachment=True, download_name="Contrato_Generado.docx")

    return render_template("form.html", fields=fields)

if __name__ == "__main__":
    app.run(debug=True)
